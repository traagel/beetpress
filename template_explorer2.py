from pptx import Presentation

# Load the existing presentation
prs = Presentation("taltech_temp1.pptx")

# Get slide number 2 (index 1 since it's 0-based)
slide_to_extract = prs.slides[1]

# Create a new presentation for the template
template_prs = Presentation()

# Add a slide to the new presentation using the layout of the extracted slide
slide_layout = slide_to_extract.slide_layout
new_slide = template_prs.slides.add_slide(slide_layout)

# Save the new presentation as a template
template_prs.save("temp.pptx")

