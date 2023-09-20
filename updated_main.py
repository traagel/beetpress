
from pptx import Presentation
from pptx.util import Inches

# Initialize a presentation object with a template
prs = Presentation("TEMPLATE.pptx")

# Get the slide layout from the template slide
template_slide_layout = prs.slides[0].slide_layout



# Read data from a text file
with open("/slide_content.txt", "r") as f:
    lines = f.readlines()

# Loop through each line in the file to create slides
for line in lines:
    slide = prs.slides.add_slide(template_slide_layout)

    # Check if the slide has a title and content placeholder
    if len(slide.placeholders) > 1:
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set the title and content for each slide
        title.text = slide_title
        content.text = line.strip()
    else:
        print(f"Skipping slide due to insufficient placeholders: {line.strip()}")

# Save the presentation
prs.save("lecture_slides.pptx")
