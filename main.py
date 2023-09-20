from pptx import Presentation

def parse_slide_content(slide_text):
    lines = slide_text.strip().split('\n')
    title = None
    bullet_points = []

    for line in lines:
        stripped_line = line.strip()

        if stripped_line.startswith('Title:'):
            title = stripped_line.split('Title:')[1].strip()
        elif stripped_line.startswith('-'):
            depth = len(stripped_line) - len(stripped_line.lstrip('-'))
            text = stripped_line[depth:].strip()
            bullet_points.append((depth, text))

    return title, bullet_points

prs = Presentation('./template/lecture_slides_tmp_v3.pptx')


xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
xml_slides.remove(slides[0])  # Delete the first slide because it's a template

content_slide_layout = prs.slide_layouts[1]

with open("./loeng3-slaidid-est.txt", "r") as f:
    content = f.read()

slides = content.split('==SLIDE==')

for slide_text in slides:
    slide_text = slide_text.strip()
    if not slide_text:
        continue

    title_text, bullet_points = parse_slide_content(slide_text)
    slide = prs.slides.add_slide(content_slide_layout)

    if title_text:
        slide.shapes.title.text = title_text

    if bullet_points:
        tf = slide.placeholders[1].text_frame
        tf.clear()  # Clear existing text

        # Handle the first bullet point
        depth, text = bullet_points.pop(0)
        tf.text = text

        # Handle the remaining bullet points
        for depth, text in bullet_points:
            p = tf.add_paragraph()
            p.text = text
            p.level = depth - 1

prs.save("lecture_slides.pptx")
